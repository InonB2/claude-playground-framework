#!/usr/bin/env python3
"""
CV Scale Pipeline — generate tailored CVs from the V4 master template.

FLOW:
  1. Andy/Claude receives a JD (URL or text file or pasted text).
  2. Andy analyses the JD, writes a tailored 3-5 sentence summary for Inon.
  3. Run: python scripts/cv_scale.py generate --role "PM" --company "Acme" \
              --location "Tel Aviv" [--req-id 1234] [--jd jd.txt] [--jd-url URL]
     → Creates: output/cv_archive/<folder>/v1_Inon_Baasov_CV_<role>.pptx
     → Also saves JD.txt into the folder.
  4. Inon reviews and approves the PPTX.
  5. Run: python scripts/cv_scale.py topdf output/cv_archive/<folder>/v1_...pptx
     → Creates a PDF alongside the PPTX.

TEMPLATE: output/cv_archive/6486_TrainingPM_Elbit_Netanya/v4_Inon_Baasov_CV_TrainingPM.pptx
"""

import argparse
import re
import shutil
import sys
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE_DIR  = Path(__file__).resolve().parent.parent
ARCHIVE   = BASE_DIR / "output" / "cv_archive"
TEMPLATE  = ARCHIVE / "6486_TrainingPM_Elbit_Netanya" / "v4_Inon_Baasov_CV_TrainingPM.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slugify(text: str) -> str:
    """Convert role/company to a safe folder-name fragment."""
    return re.sub(r"[^a-zA-Z0-9]+", "_", text).strip("_")


def build_folder_name(req_id: str | None, role: str, company: str, location: str) -> str:
    parts = []
    if req_id:
        parts.append(req_id)
    parts.append(slugify(role))
    parts.append(slugify(company))
    if location:
        parts.append(slugify(location))
    return "_".join(parts)


def fetch_jd_text(jd_file: str | None, jd_url: str | None, jd_text: str | None) -> str:
    """Return JD as a plain string from whichever source is given."""
    if jd_file:
        return Path(jd_file).read_text(encoding="utf-8")
    if jd_url:
        req = urllib.request.Request(jd_url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
        # Strip HTML tags (simple)
        return re.sub(r"<[^>]+>", " ", raw)
    if jd_text:
        return jd_text
    return ""


def _all_text_frames(shapes):
    """Yield all text frames, recursively descending into groups."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _all_text_frames(shape.shapes)
        elif shape.has_text_frame:
            yield shape.text_frame


def update_summary(prs: Presentation, new_summary: str) -> bool:
    """Replace the professional summary body text in the first slide (searches groups too)."""
    slide = prs.slides[0]
    for tf in _all_text_frames(slide.shapes):
        text = tf.text.strip()
        if (
            (text.startswith("Product leader") or text.startswith("Product Leader")
             or text.startswith("Principal Product Owner"))
            and len(text) > 100
        ):
            first_para = tf.paragraphs[0]
            first_run  = next((r for p in tf.paragraphs for r in p.runs), None)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            if first_run:
                first_run.text = new_summary
            else:
                r = first_para.add_run()
                r.text = new_summary
                r.font.size = Pt(11)
            return True
    return False


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------
def cmd_generate(args):
    if not TEMPLATE.exists():
        sys.exit(f"ERROR: template not found at {TEMPLATE}")

    folder_name = build_folder_name(args.req_id, args.role, args.company, args.location or "")
    out_dir     = ARCHIVE / folder_name
    out_dir.mkdir(parents=True, exist_ok=True)

    # Save JD
    jd_text = fetch_jd_text(args.jd, args.jd_url, args.jd_text)
    if jd_text:
        (out_dir / "JD.txt").write_text(jd_text, encoding="utf-8")
        print(f"  JD saved → {out_dir / 'JD.txt'}")

    # Copy template
    role_slug = slugify(args.role)
    out_file  = out_dir / f"v1_Inon_Baasov_CV_{role_slug}.pptx"
    shutil.copy2(TEMPLATE, out_file)
    print(f"  Template copied → {out_file.relative_to(BASE_DIR)}")

    # Apply custom summary if provided
    if args.summary:
        prs = Presentation(str(out_file))
        ok  = update_summary(prs, args.summary)
        prs.save(str(out_file))
        if ok:
            print("  Summary updated ✓")
        else:
            print("  WARN: summary placeholder not found — PPTX saved without summary change")
    else:
        print("  NOTE: no --summary provided. Open the PPTX and update the Professional")
        print("        Summary section manually, or re-run with --summary 'Your text here'.")

    print(f"\nReview: {out_file}")
    print("When approved, run:")
    print(f"  python scripts/cv_scale.py topdf \"{out_file}\"")


def cmd_topdf(args):
    """Convert an approved PPTX to PDF using PowerPoint COM."""
    pptx_path = Path(args.pptx).resolve()
    if not pptx_path.exists():
        sys.exit(f"ERROR: file not found: {pptx_path}")

    pdf_path = pptx_path.with_suffix(".pdf")

    try:
        import comtypes.client
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
        ppt_app.Visible = 1
        deck = ppt_app.Presentations.Open(str(pptx_path))
        deck.SaveAs(str(pdf_path), 32)   # 32 = ppSaveAsPDF
        deck.Close()
        ppt_app.Quit()
        print(f"PDF saved → {pdf_path}")
    except Exception as e:
        sys.exit(f"ERROR converting to PDF: {e}\nMake sure Microsoft PowerPoint is installed.")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    p = argparse.ArgumentParser(description="CV Scale Pipeline")
    sub = p.add_subparsers(dest="cmd", required=True)

    # --- generate ---
    g = sub.add_parser("generate", help="Generate a tailored CV from the master template")
    g.add_argument("--role",     required=True, help="Job role title, e.g. 'Product Manager'")
    g.add_argument("--company",  required=True, help="Company name")
    g.add_argument("--location", default="",    help="City / country")
    g.add_argument("--req-id",   default="",    dest="req_id", help="Job requisition ID")
    g.add_argument("--jd",       default=None,  help="Path to JD text file")
    g.add_argument("--jd-url",   default=None,  help="URL of the job posting")
    g.add_argument("--jd-text",  default=None,  help="Raw JD text (inline)")
    g.add_argument("--summary",  default=None,
                   help="Tailored professional summary (3-5 sentences) to inject into the CV")

    # --- topdf ---
    t = sub.add_parser("topdf", help="Convert an approved PPTX to PDF")
    t.add_argument("pptx", help="Path to the approved PPTX file")

    args = p.parse_args()
    if args.cmd == "generate":
        cmd_generate(args)
    elif args.cmd == "topdf":
        cmd_topdf(args)


if __name__ == "__main__":
    main()
