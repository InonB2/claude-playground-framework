"""
BuildARPro — Pitch Deck Generator
Generates: D:/Claude Playground/owner_inbox/BuildARPro_PitchDeck.pptx
12 slides, 16:9 widescreen, dark navy + cyan design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Colors ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x0E, 0x1A)
CYAN        = RGBColor(0x00, 0xB4, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xBE, 0xC5)
DARK_GRAY   = RGBColor(0x1A, 0x23, 0x35)   # card/panel bg
MID_GRAY    = RGBColor(0x2A, 0x35, 0x50)   # separator line

OUTPUT_PATH = r"D:\Claude Playground\owner_inbox\BuildARPro_PitchDeck.pptx"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank


def add_bg(slide):
    """Fill slide background with dark navy."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb


def add_multiline(slide, lines, left, top, width, height,
                  base_size=16, line_gap=None):
    """
    lines = list of (text, font_size, bold, color, align)
    Stacks them in a single textbox with paragraph breaks.
    """
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, fs, bold, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_gap:
            p.space_before = Pt(line_gap)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Arial"
    return txb


def add_slide_number(slide, num, total=12):
    """Small slide number bottom-right."""
    add_textbox(slide, f"{num} / {total}",
                left=Inches(12.2), top=Inches(7.1),
                width=Inches(1.0), height=Inches(0.3),
                font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def add_cyan_title(slide, title_text, top=Inches(0.35), left=Inches(0.6),
                   width=Inches(12.0), font_size=34):
    add_textbox(slide, title_text,
                left=left, top=top, width=width, height=Inches(0.6),
                font_size=font_size, bold=True, color=CYAN)


def add_separator(slide, top=Inches(1.05)):
    """Thin cyan horizontal line under the title."""
    add_rect(slide,
             left=Inches(0.6), top=top,
             width=Inches(12.13), height=Pt(2),
             fill_color=CYAN)


def bullet_block(slide, bullets, left, top, width, height,
                 font_size=16, bullet_char="•  "):
    """Renders a list of bullet strings as stacked paragraphs."""
    lines = []
    for i, b in enumerate(bullets):
        gap = 6 if i > 0 else 0
        lines.append((f"{bullet_char}{b}", font_size, False, LIGHT_GRAY, PP_ALIGN.LEFT))
    add_multiline(slide, lines, left, top, width, height, line_gap=6)


def stat_card(slide, label, value, left, top, w=Inches(2.8), h=Inches(1.4)):
    """Dark card with a big cyan value and white label."""
    add_rect(slide, left, top, w, h, fill_color=DARK_GRAY)
    add_textbox(slide, value,
                left=left + Inches(0.15), top=top + Inches(0.1),
                width=w - Inches(0.3), height=Inches(0.65),
                font_size=28, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                left=left + Inches(0.1), top=top + Inches(0.75),
                width=w - Inches(0.2), height=Inches(0.55),
                font_size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Individual slide builders ────────────────────────────────────────────────

def slide_01_cover(prs):
    """Cover slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=CYAN)

    # Company name — large
    add_textbox(slide, "BuildARPro",
                left=Inches(1.5), top=Inches(1.6),
                width=Inches(10.0), height=Inches(1.4),
                font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, "Build Like a Pro. Every Time.",
                left=Inches(1.5), top=Inches(3.1),
                width=Inches(10.0), height=Inches(0.7),
                font_size=28, bold=False, color=CYAN, align=PP_ALIGN.CENTER)

    # Sub-tagline
    add_textbox(slide, "Your IKEA paper manual becomes a live 3D AR guide.",
                left=Inches(2.0), top=Inches(3.9),
                width=Inches(9.0), height=Inches(0.5),
                font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Bottom bar
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), fill_color=DARK_GRAY)
    add_textbox(slide, "Confidential — Seed Stage  |  Inon Baasov, Founder  |  2026",
                left=Inches(0.5), top=Inches(7.05),
                width=Inches(12.0), height=Inches(0.38),
                font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Bottom accent bar
    add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), fill_color=CYAN)

    add_slide_number(slide, 1)
    return slide


def slide_02_problem(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "DIY Instructions Are Broken")
    add_separator(slide)
    add_slide_number(slide, 2)

    bullets = [
        "600 million people globally attempt DIY projects annually — most experience failure, frustration, or costly rework",
        "The state of the art: YouTube videos you pause with dirty hands; IKEA diagrams designed to confuse",
        "No tool exists that provides real-time, spatially-aware, step-by-step guidance in your physical space",
        "$80 billion wasted annually on materials, abandoned projects, and unnecessary contractor hires",
        "Paper manuals are indecipherable. Videos are passive. Text tutorials lack spatial context.",
    ]
    bullet_block(slide, bullets,
                 left=Inches(0.7), top=Inches(1.3),
                 width=Inches(11.8), height=Inches(5.5),
                 font_size=18)

    # Big stat
    stat_card(slide, "People fail their DIY projects every year", "600M",
              left=Inches(9.5), top=Inches(1.4))
    stat_card(slide, "Wasted on failed DIY annually", "$80B",
              left=Inches(9.5), top=Inches(3.0))

    return slide


def slide_03_solution(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "AR Guidance That Meets You Where You Build")
    add_separator(slide)
    add_slide_number(slide, 3)

    # Vision statement
    add_textbox(slide,
                "Point your phone at a paper manual. A live 3D AR guide appears.",
                left=Inches(0.7), top=Inches(1.3),
                width=Inches(11.8), height=Inches(0.6),
                font_size=22, bold=True, color=WHITE)

    # How it works summary
    steps = [
        ("1", "Scan", "Point your phone camera at a physical paper manual or instruction sheet"),
        ("2", "Interpret", "OCR + Computer Vision + 3D model generation reads and maps the instructions"),
        ("3", "Overlay", "A rotatable 3D model appears anchored in your real-world space via AR"),
        ("4", "Guide", "Real-time arrows, highlights, and step-by-step AR cues walk you through every step"),
    ]

    for i, (num, heading, desc) in enumerate(steps):
        col_x = Inches(0.5 + i * 3.2)
        add_rect(slide, col_x, Inches(2.1), Inches(3.0), Inches(4.5), fill_color=DARK_GRAY)
        add_textbox(slide, num,
                    left=col_x + Inches(0.1), top=Inches(2.2),
                    width=Inches(0.5), height=Inches(0.5),
                    font_size=28, bold=True, color=CYAN)
        add_textbox(slide, heading,
                    left=col_x + Inches(0.1), top=Inches(2.75),
                    width=Inches(2.8), height=Inches(0.45),
                    font_size=16, bold=True, color=WHITE)
        add_textbox(slide, desc,
                    left=col_x + Inches(0.1), top=Inches(3.25),
                    width=Inches(2.8), height=Inches(1.2),
                    font_size=13, color=LIGHT_GRAY, wrap=True)

    add_textbox(slide,
                '"Your IKEA paper manual becomes a live 3D AR guide."',
                left=Inches(1.5), top=Inches(6.7),
                width=Inches(10.0), height=Inches(0.5),
                font_size=16, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

    return slide


def slide_04_how_it_works(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "How It Works — Four Steps to Pro Results")
    add_separator(slide)
    add_slide_number(slide, 4)

    steps = [
        ("01", "Choose Your Project",
         "Browse 500+ verified DIY projects by category and skill level. Community ratings guide your choice."),
        ("02", "Scan the Manual",
         "Point your camera at any paper instruction sheet. OCR + Computer Vision + 3D interpretation activates instantly."),
        ("03", "Follow the AR Guide",
         "A rotatable 3D model overlays your real world. Step-by-step arrows, highlights, and explanations guide every move. Hands stay on the project — not your screen."),
        ("04", "Complete and Share",
         "Log your build, earn your badge, and inspire the next builder in the community feed."),
    ]

    for i, (num, heading, desc) in enumerate(steps):
        row_y = Inches(1.5 + i * 1.35)
        # Number circle area
        add_rect(slide, Inches(0.5), row_y, Inches(0.9), Inches(1.1), fill_color=DARK_GRAY)
        add_textbox(slide, num,
                    left=Inches(0.5), top=row_y + Inches(0.2),
                    width=Inches(0.9), height=Inches(0.6),
                    font_size=24, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_textbox(slide, heading,
                    left=Inches(1.55), top=row_y + Inches(0.05),
                    width=Inches(3.5), height=Inches(0.45),
                    font_size=17, bold=True, color=WHITE)
        add_textbox(slide, desc,
                    left=Inches(1.55), top=row_y + Inches(0.5),
                    width=Inches(11.0), height=Inches(0.75),
                    font_size=14, color=LIGHT_GRAY, wrap=True)

    return slide


def slide_05_market(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Three Massive Markets, One Platform")
    add_separator(slide)
    add_slide_number(slide, 5)

    # Three market cards
    markets = [
        ("$800B", "Global DIY\nHome Improvement\n(2025, growing to $1.4T by 2032)"),
        ("$30.6B", "Mobile AR Market\n(2025, growing to $113.6B by 2030\nCAGR 31.3%)"),
        ("$205B", "Creator Economy\n(2024, CAGR 23.3%\nDIY content is a primary pillar)"),
    ]
    for i, (value, label) in enumerate(markets):
        x = Inches(0.5 + i * 4.2)
        stat_card(slide, label, value,
                  left=x, top=Inches(1.4), w=Inches(3.9), h=Inches(2.5))

    # TAM row
    add_textbox(slide, "Target Addressable Market — AR-Guided DIY Tools",
                left=Inches(0.6), top=Inches(4.2),
                width=Inches(11.8), height=Inches(0.45),
                font_size=18, bold=True, color=WHITE)

    tam_stats = [
        ("$4–10B", "Near-term TAM"),
        ("$20B+", "Long-term TAM at maturity"),
        ("100M", "Potential users\n(5% of 2B AR users, DIY-engaged)"),
        ("31.3%", "Mobile AR CAGR\n2025–2030"),
    ]
    for i, (value, label) in enumerate(tam_stats):
        x = Inches(0.5 + i * 3.2)
        stat_card(slide, label, value,
                  left=x, top=Inches(4.75), w=Inches(3.0), h=Inches(1.5))

    return slide


def slide_06_competitive(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "We Start Where Everyone Else Stops")
    add_separator(slide)
    add_slide_number(slide, 6)

    add_textbox(slide, "No competitor combines AR spatial overlay + step-by-step DIY guidance.",
                left=Inches(0.7), top=Inches(1.3),
                width=Inches(11.8), height=Inches(0.45),
                font_size=18, bold=True, color=WHITE)

    # Comparison table header
    headers = ["Product", "AR Overlay", "Step-by-Step Guidance", "Works on Any Project", "AI Assistant"]
    col_widths = [Inches(2.5), Inches(2.0), Inches(3.0), Inches(2.5), Inches(2.0)]
    col_starts = [Inches(0.5), Inches(3.0), Inches(5.0), Inches(8.0), Inches(10.5)]

    row_y = Inches(1.9)
    add_rect(slide, Inches(0.4), row_y, Inches(12.5), Inches(0.5), fill_color=MID_GRAY)
    for j, h in enumerate(headers):
        add_textbox(slide, h,
                    left=col_starts[j], top=row_y + Inches(0.05),
                    width=col_widths[j], height=Inches(0.4),
                    font_size=13, bold=True, color=CYAN)

    rows = [
        ("IKEA Place",       "Yes",  "No",      "No (IKEA only)", "No"),
        ("Wayfair AR",       "Yes",  "No",      "No (Wayfair only)", "No"),
        ("YouTube DIY",      "No",   "Partial", "Yes",            "No"),
        ("Instructables",    "No",   "Yes",     "Yes",            "No"),
        ("Google Lens",      "No",   "No",      "Limited",        "Partial"),
        ("BuildARPro",      "Yes",  "Yes",     "Yes",            "Yes"),
    ]

    for i, row in enumerate(rows):
        ry = Inches(2.45 + i * 0.65)
        bg = DARK_GRAY if i % 2 == 0 else NAVY
        if row[0] == "BuildARPro":
            bg = RGBColor(0x00, 0x2A, 0x3A)  # highlighted row
        add_rect(slide, Inches(0.4), ry, Inches(12.5), Inches(0.6), fill_color=bg)
        for j, cell in enumerate(row):
            color = CYAN if row[0] == "BuildARPro" else (
                WHITE if j == 0 else (LIGHT_GRAY if cell not in ("No", "No (IKEA only)", "No (Wayfair only)") else RGBColor(0x66, 0x77, 0x88))
            )
            add_textbox(slide, cell,
                        left=col_starts[j], top=ry + Inches(0.1),
                        width=col_widths[j], height=Inches(0.4),
                        font_size=13, color=color,
                        bold=(row[0] == "BuildARPro"))

    return slide


def slide_07_business_model(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Three Revenue Streams, One Network")
    add_separator(slide)
    add_slide_number(slide, 7)

    streams = [
        ("Consumer Pro\nSubscription",
         "$7.99/month\n$59.99/year",
         "Unlimited projects, full AI assistant,\nvoice guidance, offline downloads,\nadvanced AR features, early access."),
        ("Creator\nMarketplace",
         "Revenue Share",
         "Verified creators publish projects and\nearn % of Pro subscriber completions.\nPlatform takes 30% — like App Store."),
        ("B2B\nWhite-Label",
         "$5K–$50K/mo",
         "Furniture brands and hardware retailers\nlicense BuildARPro to replace paper\nassembly manuals. Enterprise licensing."),
    ]

    for i, (title, value, desc) in enumerate(streams):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(1.4), Inches(3.9), Inches(4.2), fill_color=DARK_GRAY)
        add_textbox(slide, title,
                    left=x + Inches(0.15), top=Inches(1.5),
                    width=Inches(3.6), height=Inches(0.75),
                    font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, value,
                    left=x + Inches(0.1), top=Inches(2.35),
                    width=Inches(3.7), height=Inches(0.65),
                    font_size=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    left=x + Inches(0.15), top=Inches(3.1),
                    width=Inches(3.6), height=Inches(1.35),
                    font_size=13, color=LIGHT_GRAY, wrap=True, align=PP_ALIGN.CENTER)

    # Projections
    add_textbox(slide, "Revenue Projections",
                left=Inches(0.6), top=Inches(5.85),
                width=Inches(4.0), height=Inches(0.4),
                font_size=14, bold=True, color=WHITE)
    add_textbox(slide,
                "Year 1: 50,000 Pro subscribers = $4.8M ARR    |    Year 3: 500,000 subscribers + B2B = $50M+ ARR",
                left=Inches(0.6), top=Inches(6.3),
                width=Inches(12.2), height=Inches(0.4),
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    return slide


def slide_08_technology(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Technology — Built for the Real World")
    add_separator(slide)
    add_slide_number(slide, 8)

    tech_layers = [
        ("OCR Engine",
         "Optical Character Recognition scans paper manuals in real-time, extracting step data, part numbers, and spatial relationships from physical instruction sheets."),
        ("Computer Vision + 3D Interpretation",
         "CV pipeline interprets 2D manual diagrams as 3D spatial models — identifying components, their relationships, and assembly sequences automatically."),
        ("ARKit (iOS) / ARCore (Android)",
         "Native AR frameworks enable precise plane detection, LiDAR support on iPhone 12 Pro+, world tracking, and spatial anchoring of 3D overlays in real space."),
        ("Supabase Backend",
         "PostgreSQL database, Auth, Realtime, Edge Functions, and Storage — handles project library, user accounts, creator content, community, and AI assistant API calls."),
    ]

    for i, (title, desc) in enumerate(tech_layers):
        ry = Inches(1.4 + i * 1.4)
        # Left accent bar
        add_rect(slide, Inches(0.5), ry, Inches(0.07), Inches(1.1), fill_color=CYAN)
        add_textbox(slide, title,
                    left=Inches(0.75), top=ry,
                    width=Inches(3.5), height=Inches(0.45),
                    font_size=15, bold=True, color=WHITE)
        add_textbox(slide, desc,
                    left=Inches(0.75), top=ry + Inches(0.45),
                    width=Inches(11.8), height=Inches(0.75),
                    font_size=13, color=LIGHT_GRAY, wrap=True)

    # Stack summary
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), fill_color=DARK_GRAY)
    add_textbox(slide,
                "Stack: React Native  +  ViroReact  +  ARKit/ARCore  +  OCR  +  Computer Vision  +  Supabase  +  AI (Claude/OpenAI)",
                left=Inches(0.6), top=Inches(7.03),
                width=Inches(12.1), height=Inches(0.3),
                font_size=12, color=CYAN, align=PP_ALIGN.CENTER)

    return slide


def slide_09_traction(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Early Signals — Stealth Development")
    add_separator(slide)
    add_slide_number(slide, 9)

    # Status badge
    add_rect(slide, Inches(9.5), Inches(0.25), Inches(3.3), Inches(0.55), fill_color=DARK_GRAY)
    add_textbox(slide, "Status: Stealth — Early Development",
                left=Inches(9.55), top=Inches(0.3),
                width=Inches(3.2), height=Inches(0.42),
                font_size=13, bold=True, color=CYAN, align=PP_ALIGN.RIGHT)

    add_textbox(slide,
                "Israel market validation in progress. No public launch yet. Building the right product before marketing it.",
                left=Inches(0.7), top=Inches(1.35),
                width=Inches(11.8), height=Inches(0.5),
                font_size=16, color=LIGHT_GRAY)

    validation_points = [
        ("Market Validation", "Israel DIY market research complete. Target customer interviews with early testers underway. Core problem confirmed: paper manuals and passive video fail the modern DIYer."),
        ("Technical Proof of Concept", "Core AR pipeline validated: OCR scanning of paper manuals, computer vision 3D interpretation, and ARKit spatial overlay all proven in isolation. Integration in progress."),
        ("Product Vision Lock", "Corrected product vision finalized: smartphone scans a paper manual, a live 3D AR guide appears. Stealth mode — building before announcing."),
        ("Why Israel First", "Israel is an ideal seed market: high smartphone penetration, strong DIY/maker culture, tech-savvy early adopters, accessible distribution via Israeli app stores and maker communities."),
    ]

    for i, (title, desc) in enumerate(validation_points):
        ry = Inches(2.05 + i * 1.2)
        add_rect(slide, Inches(0.5), ry, Inches(0.07), Inches(1.0), fill_color=CYAN)
        add_textbox(slide, title,
                    left=Inches(0.75), top=ry,
                    width=Inches(3.0), height=Inches(0.42),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(slide, desc,
                    left=Inches(0.75), top=ry + Inches(0.42),
                    width=Inches(12.0), height=Inches(0.65),
                    font_size=13, color=LIGHT_GRAY, wrap=True)

    return slide


def slide_10_roadmap(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Roadmap — From Stealth to Scale")
    add_separator(slide)
    add_slide_number(slide, 10)

    phases = [
        ("Phase 1\nMonth 0–3\nMVP",
         [
             "Native iOS + Android app (React Native + ViroReact)",
             "OCR + CV pipeline: scan any paper manual",
             "AR step-by-step overlay with arrows and 3D guides",
             "20–30 curated seed projects across 4 categories",
             "AI assistant (context-aware per project step)",
             "App Store + Google Play launch",
         ]),
        ("Phase 2\nMonth 3–6\nCreator + Community",
         [
             "Creator Studio: author AR projects, upload 3D assets",
             "Freemium paywall: Pro at $7.99/month",
             "Creator revenue share program",
             "Community feed, social sharing, comments",
             "Voice guidance (TTS) during AR sessions",
             "Image recognition: scan a part to identify it",
         ]),
        ("Phase 3\nMonth 6–12\nScale + Intelligence",
         [
             "Full 3D assembly reconstruction (LiDAR on capable devices)",
             "Error detection: AI flags misaligned components",
             "B2B white-label channel: IKEA, Home Depot, West Elm",
             "Collaborative mode: two users, one project",
             "500+ verified projects, 6+ categories",
             "Offline mode + AR project recording",
         ]),
    ]

    for i, (title, items) in enumerate(phases):
        x = Inches(0.45 + i * 4.3)
        add_rect(slide, x, Inches(1.4), Inches(4.0), Inches(5.7), fill_color=DARK_GRAY)
        add_textbox(slide, title,
                    left=x + Inches(0.15), top=Inches(1.5),
                    width=Inches(3.7), height=Inches(0.85),
                    font_size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_rect(slide, x + Inches(0.15), Inches(2.45), Inches(3.7), Pt(1.5), fill_color=MID_GRAY)
        for j, item in enumerate(items):
            add_textbox(slide, f"•  {item}",
                        left=x + Inches(0.15), top=Inches(2.6 + j * 0.75),
                        width=Inches(3.7), height=Inches(0.7),
                        font_size=12, color=LIGHT_GRAY, wrap=True)

    return slide


def slide_11_team(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Team — Built by People Who Know the Domain")
    add_separator(slide)
    add_slide_number(slide, 11)

    # Founder block
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.2), fill_color=DARK_GRAY)

    add_textbox(slide, "Inon Baasov",
                left=Inches(0.7), top=Inches(1.5),
                width=Inches(6.0), height=Inches(0.6),
                font_size=28, bold=True, color=WHITE)
    add_textbox(slide, "Product Leader & Founder",
                left=Inches(0.7), top=Inches(2.1),
                width=Inches(6.0), height=Inches(0.4),
                font_size=18, bold=False, color=CYAN)

    founder_bullets = [
        "Technion (Israel Institute of Technology) graduate",
        "$2.5M+ product impact across prior roles",
        "Deep expertise in product strategy, user experience, and go-to-market execution",
        "Solo founder — product-led, Israel-first, stealth stage",
        "Vision: democratize professional-quality DIY guidance for every smartphone user",
    ]
    for j, b in enumerate(founder_bullets):
        add_textbox(slide, f"•  {b}",
                    left=Inches(0.7), top=Inches(2.65 + j * 0.38),
                    width=Inches(11.8), height=Inches(0.35),
                    font_size=14, color=LIGHT_GRAY)

    # Hiring plan
    add_textbox(slide, "Hiring Plan — Funded Roles",
                left=Inches(0.6), top=Inches(4.85),
                width=Inches(6.0), height=Inches(0.4),
                font_size=16, bold=True, color=WHITE)

    hiring = [
        ("AR / React Native Engineer", "Core AR pipeline, iOS + Android", "60% of seed engineering budget"),
        ("UX / Product Designer", "App and Creator Studio interfaces", "Equity + cash"),
        ("Creator Acquisition Lead", "Recruit first 50 verified creators before launch", "Commission + equity"),
    ]
    for j, (role, focus, comp) in enumerate(hiring):
        ry = Inches(5.35 + j * 0.6)
        add_textbox(slide, f"{role}  —  {focus}  ({comp})",
                    left=Inches(0.7), top=ry,
                    width=Inches(12.0), height=Inches(0.5),
                    font_size=13, color=LIGHT_GRAY)

    add_slide_number(slide, 11)
    return slide


def slide_12_ask(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide)
    add_cyan_title(slide, "Join Us in Redefining How the World Builds")
    add_separator(slide)
    add_slide_number(slide, 12)

    # The ask
    add_textbox(slide, "The Ask",
                left=Inches(0.7), top=Inches(1.3),
                width=Inches(5.0), height=Inches(0.5),
                font_size=22, bold=True, color=CYAN)

    add_textbox(slide,
                "Raising Pre-Seed / Seed capital to fund AR MVP development and initial creator + user acquisition.",
                left=Inches(0.7), top=Inches(1.85),
                width=Inches(11.8), height=Inches(0.5),
                font_size=16, color=LIGHT_GRAY)

    # Use of funds
    funds = [
        ("60%", "Engineering", "AR development, native app, OCR + CV pipeline, Supabase integration"),
        ("25%", "Creator Acquisition", "Recruit first 50 verified creators, seed content production"),
        ("15%", "Operations", "Infrastructure, App Store fees, legal, working capital"),
    ]

    for i, (pct, label, desc) in enumerate(funds):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(2.6), Inches(3.9), Inches(2.0), fill_color=DARK_GRAY)
        add_textbox(slide, pct,
                    left=x + Inches(0.1), top=Inches(2.7),
                    width=Inches(3.7), height=Inches(0.7),
                    font_size=36, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_textbox(slide, label,
                    left=x + Inches(0.1), top=Inches(3.4),
                    width=Inches(3.7), height=Inches(0.3),
                    font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    left=x + Inches(0.1), top=Inches(3.75),
                    width=Inches(3.7), height=Inches(0.7),
                    font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, wrap=True)

    # Milestones
    add_textbox(slide, "Milestones With Funding",
                left=Inches(0.6), top=Inches(4.85),
                width=Inches(5.0), height=Inches(0.4),
                font_size=16, bold=True, color=WHITE)

    milestones = [
        ("Month 3", "AR MVP live in App Store — OCR scans any paper manual, 3D AR guide appears"),
        ("Month 6", "10,000 users, Pro subscription revenue live, 50+ verified creator projects"),
        ("Month 9", "First B2B letter of intent from furniture retailer or hardware brand"),
        ("Month 12", "50,000 users, Creator Marketplace active, B2B channel generating revenue"),
    ]
    for j, (month, goal) in enumerate(milestones):
        add_textbox(slide, f"{month}: {goal}",
                    left=Inches(0.7), top=Inches(5.35 + j * 0.43),
                    width=Inches(12.0), height=Inches(0.4),
                    font_size=13, color=LIGHT_GRAY)

    # Contact
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), fill_color=DARK_GRAY)
    add_textbox(slide,
                "Inon Baasov  |  inonbaasov@gmail.com  |  promakerapp.com  |  \"Build Like a Pro. Every Time.\"",
                left=Inches(0.5), top=Inches(7.05),
                width=Inches(12.3), height=Inches(0.38),
                font_size=12, color=CYAN, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), fill_color=CYAN)

    return slide


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_cover(prs)
    print("  Slide 1: Cover")
    slide_02_problem(prs)
    print("  Slide 2: Problem")
    slide_03_solution(prs)
    print("  Slide 3: Solution")
    slide_04_how_it_works(prs)
    print("  Slide 4: How It Works")
    slide_05_market(prs)
    print("  Slide 5: Market")
    slide_06_competitive(prs)
    print("  Slide 6: Competitive")
    slide_07_business_model(prs)
    print("  Slide 7: Business Model")
    slide_08_technology(prs)
    print("  Slide 8: Technology")
    slide_09_traction(prs)
    print("  Slide 9: Traction")
    slide_10_roadmap(prs)
    print("  Slide 10: Roadmap")
    slide_11_team(prs)
    print("  Slide 11: Team")
    slide_12_ask(prs)
    print("  Slide 12: The Ask")

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    size_kb = os.path.getsize(OUTPUT_PATH) // 1024
    print(f"\nSaved: {OUTPUT_PATH}")
    print(f"Size:  {size_kb} KB ({os.path.getsize(OUTPUT_PATH):,} bytes)")
    print("Done.")


if __name__ == "__main__":
    main()
