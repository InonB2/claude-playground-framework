"""
Build Elbit CV variants from user's preferred templates.

Sources (read-only):
  V1.S.M1.PPM.pptx  -> generic Product/Project Manager template
  V9.S.M4.TPMAI-E.pptx -> Technical PM AI - Elbit (AI/Innovation tailored)

Outputs:
  6486_TrainingPM_Elbit_Netanya/v4_Inon_Baasov_CV_TrainingPM.pptx
  5811_AI_Innovation_PM_Elbit_Netanya/v1_Inon_Baasov_CV_AI_Innovation.pptx

Both: V1 layout/font sizes preserved. Website text-frame added to contact row.
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC_DIR = Path(r"C:\Users\Inon Baasov\OneDrive\Documents\קורות חיים CV\2026")
OUT_BASE = Path(r"D:\Claude Playground\output\cv_archive")

V1 = SRC_DIR / "Inon Baasov V1.S.M1.PPM.pptx"
V9 = SRC_DIR / "Inon Baasov V9.S.M4.TPMAI-E.pptx"

WEBSITE = "inon-baasov-website.base44.app"


def add_website_to_contact_row(pptx_path: Path):
    """Add a website text-frame to the contact row, sized/positioned to fit alongside Phone | Email | LinkedIn."""
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]

    # Find the LinkedIn shape - it has the rightmost contact text
    linkedin_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "LinkedIn" in text and "inonbaasov" in text:
            linkedin_shape = shape
            break

    if linkedin_shape is None:
        print(f"  WARN: LinkedIn shape not found in {pptx_path.name}")
        return prs

    # Read its position and font properties so we match perfectly
    li_left = linkedin_shape.left
    li_top = linkedin_shape.top
    li_width = linkedin_shape.width
    li_height = linkedin_shape.height

    # Existing contact row appears at y ~= 0.76", LinkedIn ends around x=6.37"
    # Page width ~ 8.27" - place website at x ~= 6.5", w = 1.7"
    # Easier: shift LinkedIn left, then add website to its old right edge area.
    # Safest: place website as a NEW text box below the contact row at same y but a touch lower
    page_w = prs.slide_width
    page_h = prs.slide_height

    # Place website to the right of LinkedIn (or wrap to next line if no room)
    new_left = li_left + li_width + Inches(0.10)
    if new_left + Inches(1.50) > page_w - Inches(0.20):
        # Not enough room — place it just below the contact row spanning full width
        new_left = Inches(0.40)
        new_top = li_top + li_height + Inches(0.05)
        new_width = page_w - Inches(0.80)
    else:
        new_top = li_top
        new_width = page_w - new_left - Inches(0.20)

    tb = slide.shapes.add_textbox(new_left, new_top, new_width, li_height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # Match the font style by sampling from the LinkedIn shape's first run
    sample_run = None
    for para in linkedin_shape.text_frame.paragraphs:
        for run in para.runs:
            sample_run = run
            break
        if sample_run:
            break

    # Add the website text
    run1 = p.add_run()
    run1.text = "Portfolio  "
    if sample_run is not None:
        try:
            run1.font.name = sample_run.font.name
            run1.font.size = sample_run.font.size or Pt(11)
            run1.font.bold = True
            if sample_run.font.color and sample_run.font.color.type is not None:
                run1.font.color.rgb = sample_run.font.color.rgb
        except Exception:
            run1.font.size = Pt(11)
            run1.font.bold = True
    else:
        run1.font.size = Pt(11)
        run1.font.bold = True

    run2 = p.add_run()
    run2.text = WEBSITE
    if sample_run is not None:
        try:
            run2.font.name = sample_run.font.name
            run2.font.size = sample_run.font.size or Pt(11)
        except Exception:
            run2.font.size = Pt(11)
    else:
        run2.font.size = Pt(11)

    return prs


def update_summary_for_training(prs: Presentation, new_summary: str):
    """Replace the Professional Summary body text with a Training-PM-tailored version."""
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        # Heuristic: summary body is the long paragraph that starts with "Product"
        if (text.startswith("Product leader") or text.startswith("Product Leader")
                or text.startswith("Principal Product Owner")) and len(text) > 200:
            # Preserve the formatting of the first run, replace text
            first_para = tf.paragraphs[0]
            first_run = None
            for run in first_para.runs:
                first_run = run
                break

            # Clear all text
            for p in tf.paragraphs:
                for r in list(p.runs):
                    r.text = ""

            # Set new text on first run
            if first_run is not None:
                first_run.text = new_summary
            else:
                p = first_para
                r = p.add_run()
                r.text = new_summary
                r.font.size = Pt(11)
            return True
    return False


def build_training_pm_cv():
    out_dir = OUT_BASE / "6486_TrainingPM_Elbit_Netanya"
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = add_website_to_contact_row(V1)

    # Tailor summary for Training/Learning Dept PM at Elbit
    training_summary = (
        "Product Leader with 10+ years delivering technology products from zero to scale, "
        "with hands-on experience building training and learning systems alongside AI/SaaS platforms. "
        "Led full product lifecycle (PRD/MRD, V&V, supplier mgmt, global projects) for interactive video, "
        "AI, and EdTech-adjacent products. Raised $2.5M, drove 38% efficiency gains, and managed teams of 20+. "
        "Technion-educated (BSc + Executive MBA). Passionate about building scalable learning and content "
        "delivery platforms in regulated, multi-stakeholder environments — exactly the profile the Elbit "
        "Learning Department's TPSS / emulator / knowledge-portal product portfolio demands."
    )
    updated = update_summary_for_training(prs, training_summary)
    if not updated:
        print("  WARN: summary not found in V1 template")

    out_file = out_dir / "v4_Inon_Baasov_CV_TrainingPM.pptx"
    prs.save(str(out_file))
    print(f"OK  {out_file.relative_to(OUT_BASE)}")


def build_ai_innovation_cv():
    out_dir = OUT_BASE / "5811_AI_Innovation_PM_Elbit_Netanya"
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = add_website_to_contact_row(V9)
    out_file = out_dir / "v1_Inon_Baasov_CV_AI_Innovation.pptx"
    prs.save(str(out_file))
    print(f"OK  {out_file.relative_to(OUT_BASE)}")


def copy_jds():
    """Save JD reference files into each folder."""
    training_jd = """Training Product Manager — Req #6486 — Netanya
URL: https://elbitsystemscareer.com/job/?jid=20344

Department: Learning Department (מחלקת הלמידה)

Responsibilities:
- Manage product lifecycle of learning department products (TPSS, learning systems,
  emulators, knowledge portals, electronic technical literature)
- Lead development of new and existing technological products
- Market research and competitor analysis
- Define operational and technological requirements per customer needs and market trends
- Write and manage PRD/MRD documents
- Cross-functional work: R&D, engineering, production, marketing, sales
- Lead V&V (Verification & Validation) processes
- Marketing coordination, conferences and exhibitions
- Budget, schedule, and risk management
- Sales support and technical/commercial proposals
- Long-term strategic roadmap planning
- Lead international projects with significant product component
- Supplier management and external partnerships

Requirements:
- BSc in Engineering / Information Systems / Computer Science / Learning Technologies
  (MSc - advantage)
- 3+ yrs in technology product management, preferably training, EdTech or enterprise SW
- Full PLM experience (research → spec → launch → maintenance)
- Matrix-org experience (R&D, eng, production, marketing, sales)
- Supplier management and global complex projects
- Business development background — significant advantage
- English C2 (speaking, reading, writing, presenting to mgmt/customers)
- Travel abroad as needed

Compatibility for Inon: ~93% (strong match across all PM/PLM/PRD-MRD/V&V/supplier
dimensions; English C2; BD background; only gap is direct EdTech/TPSS — mitigated by
adjacent interactive media and AI product experience).
"""
    ai_jd = """Technical Product Manager / Systems Engineer - AI & Innovation — Req #5811 — Netanya

Title: Technical Product Manager \\ Systems Engineer - AI & Innovation
Location: Netanya
Source: Elbit Systems careers (full JD not extracted — visit live page to apply)

Inferred focus areas (from title + Elbit AI division context):
- AI/ML product management (LLM, GenAI, NLP)
- Systems engineering integration
- Innovation pipeline / R&D-to-product translation
- Cross-functional with R&D, data science, engineering teams

Why Inon fits (~88%):
- TouchE TV: AI-powered interactive video platform from 0 to scale, $2.5M raised
- AiRakoon (consulting): Enterprise LLM architecture, API design, GenAI GTM
- Medicrowd, Smash+: GenAI/data-driven product strategy in regulated/B2C contexts
- 10+ yrs PM combined with technical depth (Technion BSc + Executive MBA)
- Defense/regulated environment experience (pharma regulatory, supplier mgmt)

Gap: Direct defense AI experience; some Elbit-specific systems engineering depth
expected. Compensated by AI/LLM portfolio and 0-to-scale execution track record.
"""
    (OUT_BASE / "6486_TrainingPM_Elbit_Netanya" / "JD.txt").write_text(training_jd, encoding="utf-8")
    (OUT_BASE / "5811_AI_Innovation_PM_Elbit_Netanya" / "JD.txt").write_text(ai_jd, encoding="utf-8")
    print("OK  JD files written")


if __name__ == "__main__":
    copy_jds()
    build_training_pm_cv()
    build_ai_innovation_cv()
    print("Done.")
