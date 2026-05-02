"""
BuildARPro — PPTX Styling Script
=================================
Reads the existing BuildARPro pitch deck and applies the brand design system:
  - Dark slide backgrounds (#0D0D0D)
  - White bold titles
  - Light gray body text (#E5E5E5)
  - Orange (#F97316) for "VIDEOS" and "BETTER" in the catchphrase
  - Matching orange for all text-run mentions of those power words

Usage:
    python buildarpro_pptx_style.py

Input:  D:\\Claude Playground\\owner_inbox\\BuildARPro_PitchDeck.pptx
Output: D:\\Claude Playground\\scratchpad\\BuildARPro_PitchDeck_Styled.pptx

Requirements:
    pip install python-pptx

Color reference (adjust constants below to change brand colors):
    DARK_BG      = #0D0D0D  — slide background
    ORANGE       = #F97316  — accent / power word color
    WHITE        = #FFFFFF  — title text
    LIGHT_GRAY   = #E5E5E5  — body text
    POWER_WORDS  = {"VIDEOS", "BETTER"}  — words that receive orange + +4pt
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy
import os

# ─────────────────────────────────────────────
# BRAND COLOR CONSTANTS — edit these to change brand colors
# ─────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x0D, 0x0D)   # Dark background for all slides
ORANGE       = RGBColor(0xF9, 0x73, 0x16)   # Primary orange accent (#F97316)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # Title and highlighted text
LIGHT_GRAY   = RGBColor(0xE5, 0xE5, 0xE5)  # Body text (#E5E5E5)
MID_GRAY     = RGBColor(0xA3, 0xA3, 0xA3)  # Secondary / caption text

# Power words that get orange color + font-size boost
POWER_WORDS  = {"VIDEOS", "BETTER"}
POWER_SIZE_BOOST_PT = 4  # How many points larger the power words become

# File paths
INPUT_PATH  = r"D:\Claude Playground\owner_inbox\BuildARPro_PitchDeck.pptx"
OUTPUT_PATH = r"D:\Claude Playground\scratchpad\BuildARPro_PitchDeck_Styled.pptx"


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def set_slide_background(slide, color: RGBColor):
    """
    Set a solid-fill background color on a slide.
    Works by manipulating the slide's XML background element directly,
    which is more reliable than the python-pptx background API.
    """
    # Get or create background element
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def is_title_placeholder(placeholder):
    """Return True if the placeholder is a title type."""
    # Placeholder type 15 = center title, 13 = title
    title_types = {0, 13, 15}  # PP_PLACEHOLDER.TITLE, CENTER_TITLE, etc.
    try:
        return placeholder.placeholder_format.idx == 0 or \
               placeholder.placeholder_format.type in title_types
    except Exception:
        return False


def style_title_run(run):
    """Apply white bold styling to a title text run."""
    run.font.bold = True
    run.font.color.rgb = WHITE


def style_body_run(run):
    """Apply light gray styling to a body text run."""
    run.font.color.rgb = LIGHT_GRAY


def style_power_word_run(run, base_size_pt):
    """
    Apply orange color and +4pt size boost to a power word run.
    base_size_pt: the font size in points before the boost (float or None).
    """
    run.font.color.rgb = ORANGE
    run.font.bold = True
    if base_size_pt is not None:
        run.font.size = Pt(base_size_pt + POWER_SIZE_BOOST_PT)
    # If no base size, we still apply the color — the size stays inherited


def split_run_on_power_words(paragraph, run_idx: int) -> int:
    """
    If a run contains a power word substring, split it into up to 3 runs:
      [before] [POWER_WORD] [after]
    Returns the new index to continue iteration from (after the inserted runs).

    This mutates the paragraph's runs list via XML manipulation.
    """
    run = paragraph.runs[run_idx]
    text = run.text
    text_upper = text.upper()

    # Find which power word appears first
    found_word = None
    found_pos = len(text)
    for word in POWER_WORDS:
        pos = text_upper.find(word)
        if pos != -1 and pos < found_pos:
            found_word = word
            found_pos = pos

    if found_word is None:
        return run_idx + 1  # No power word, move to next

    # Determine actual casing in original text
    actual_word = text[found_pos: found_pos + len(found_word)]
    before_text = text[:found_pos]
    after_text  = text[found_pos + len(found_word):]

    # Base font size for boost calculation
    base_size_pt = run.font.size.pt if run.font.size else None

    # Get the underlying XML <a:r> element
    r_elem = run._r

    # --- Build the power-word run (clone formatting from original) ---
    power_r = copy.deepcopy(r_elem)
    # Set text
    power_r.find(qn('a:t')).text = actual_word
    # Apply orange + bold + size boost to the cloned run's rPr
    rPr = power_r.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.SubElement(power_r, qn('a:rPr'))
        power_r.insert(0, rPr)
    # Remove existing solid fill if any
    for existing in rPr.findall(qn('a:solidFill')):
        rPr.remove(existing)
    # Add orange solid fill
    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', 'F97316')
    # Bold
    rPr.set('b', '1')
    # Font size
    if base_size_pt is not None:
        new_sz = int((base_size_pt + POWER_SIZE_BOOST_PT) * 100)
        rPr.set('sz', str(new_sz))

    # --- Build after run (clone) ---
    after_r = copy.deepcopy(r_elem)
    after_r.find(qn('a:t')).text = after_text

    # --- Modify original run to be "before" text ---
    r_elem.find(qn('a:t')).text = before_text

    # Insert power_r and after_r after the current run in the XML tree
    parent = r_elem.getparent()
    idx = list(parent).index(r_elem)
    parent.insert(idx + 1, power_r)
    parent.insert(idx + 2, after_r)

    # Return index of the after_r run so we continue scanning from there
    return run_idx + 2


def style_text_frame(tf, is_title: bool):
    """
    Apply brand styling to all runs in a text frame.
    - Titles: white, bold
    - Body: light gray, with power words highlighted in orange
    """
    for para in tf.paragraphs:
        run_idx = 0
        while run_idx < len(para.runs):
            run = para.runs[run_idx]
            text_upper = run.text.upper()

            # Check if this run contains a power word
            has_power = any(pw in text_upper for pw in POWER_WORDS)

            if has_power and not is_title:
                # Split and style
                run_idx = split_run_on_power_words(para, run_idx)
                # Style the power word run (now at run_idx - 1)
                # (split_run_on_power_words already set the XML attrs directly,
                #  but we re-check in case of partial matches)
            else:
                if is_title:
                    style_title_run(run)
                else:
                    # Check exact word match after splitting
                    if run.text.upper().strip() in POWER_WORDS:
                        base_size = run.font.size.pt if run.font.size else None
                        style_power_word_run(run, base_size)
                    else:
                        style_body_run(run)
                run_idx += 1


def process_shape(shape, is_title: bool):
    """Style a single shape's text frame if it has one."""
    if not shape.has_text_frame:
        return
    style_text_frame(shape.text_frame, is_title)


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────

def main():
    print(f"Loading: {INPUT_PATH}")
    if not os.path.exists(INPUT_PATH):
        raise FileNotFoundError(f"Input file not found: {INPUT_PATH}")

    prs = Presentation(INPUT_PATH)
    slide_count = len(prs.slides)
    print(f"Found {slide_count} slides. Applying BuildARPro brand styling...")

    for slide_num, slide in enumerate(prs.slides, start=1):
        print(f"  Slide {slide_num}/{slide_count}", end=" — ")

        # 1. Set dark background
        set_slide_background(slide, DARK_BG)

        # 2. Style each shape
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Detect if this shape is a title placeholder
            title_shape = False
            if shape.is_placeholder:
                try:
                    ph_idx = shape.placeholder_format.idx
                    # idx 0 = title, 12 = center title in some layouts
                    title_shape = ph_idx in (0, 12)
                except Exception:
                    pass

            process_shape(shape, is_title=title_shape)

        print("styled")

    # Save output
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"\nDone. Saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
